from __future__ import annotations

import argparse
import hashlib
import json
import os
from pathlib import Path
from typing import Any


SUPPORTED_TEXT = {".csv", ".html", ".md", ".mmd", ".txt"}
SUPPORTED_OFFICE = {".docx", ".pdf", ".pptx", ".xlsx"}
MAX_TEXT_CHARS = 250_000


def sha256_bytes(value: bytes) -> str:
    return hashlib.sha256(value).hexdigest()


def clean_text(value: str) -> str:
    lines = [line.rstrip() for line in value.replace("\x00", "").splitlines()]
    compact: list[str] = []
    blank = False
    for line in lines:
        if line.strip():
            compact.append(line)
            blank = False
        elif not blank:
            compact.append("")
            blank = True
    return "\n".join(compact).strip()[:MAX_TEXT_CHARS]


def read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                parts.append(" | ".join(values))
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for index, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"[Page {index + 1}]\n{text}")
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides):
        values: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                values.append(text.strip())
        if values:
            parts.append(f"[Slide {index + 1}]\n" + "\n".join(values))
    return "\n\n".join(parts)


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                while values and not values[-1]:
                    values.pop()
                if any(values):
                    parts.append(" | ".join(values))
    finally:
        workbook.close()
    return "\n".join(parts)


def extract(path: Path) -> tuple[str, str | None]:
    extension = path.suffix.lower()
    try:
        if extension in SUPPORTED_TEXT:
            return clean_text(read_text(path)), None
        if extension == ".docx":
            return clean_text(extract_docx(path)), None
        if extension == ".pdf":
            return clean_text(extract_pdf(path)), None
        if extension == ".pptx":
            return clean_text(extract_pptx(path)), None
        if extension == ".xlsx":
            return clean_text(extract_xlsx(path)), None
        return "", "unsupported"
    except Exception as error:  # individual artifacts report independently
        return "", f"{type(error).__name__}: {error}"


def record_for(root: Path, path: Path) -> dict[str, Any]:
    raw = path.read_bytes()
    text, error = extract(path)
    details = path.stat()
    return {
        "path": path.relative_to(root).as_posix(),
        "name": path.name,
        "extension": path.suffix.lower().lstrip("."),
        "bytes": details.st_size,
        "modifiedAt": details.st_mtime,
        "sha256": sha256_bytes(raw),
        "content": text,
        "contentChars": len(text),
        "state": "extracted" if text else ("error" if error and error != "unsupported" else "metadata-only"),
        "error": error,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--root", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    root = Path(args.root).resolve()
    output = Path(args.output).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)

    allowed = SUPPORTED_TEXT | SUPPORTED_OFFICE
    files = sorted(
        path
        for path in root.rglob("*")
        if path.is_file() and not path.name.startswith("~$") and path.suffix.lower() in allowed
    )

    counts = {"extracted": 0, "metadata-only": 0, "error": 0}
    with output.open("w", encoding="utf-8", newline="\n") as handle:
        for path in files:
            record = record_for(root, path)
            counts[record["state"]] = counts.get(record["state"], 0) + 1
            handle.write(json.dumps(record, ensure_ascii=False) + "\n")

    print(
        json.dumps(
            {
                "ok": True,
                "root": str(root),
                "output": str(output),
                "files": len(files),
                "counts": counts,
            }
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
